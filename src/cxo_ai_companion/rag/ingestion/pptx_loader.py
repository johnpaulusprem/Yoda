"""PPTX (PowerPoint) document loader using python-pptx."""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from typing import Any

from cxo_ai_companion.rag.ingestion.base_loader import (
    DocumentLoader,
    DocumentMetadata,
    LoadedDocument,
    LoaderConfig,
    LoadMode,
)

logger = logging.getLogger(__name__)


@dataclass
class PPTXConfig(LoaderConfig):
    """PPTX-specific loader configuration."""

    extract_notes: bool = True
    extract_tables: bool = True


class PPTXLoader(DocumentLoader):
    """Load text content from PPTX presentations via python-pptx."""

    def __init__(self, config: PPTXConfig | None = None) -> None:
        super().__init__(config or PPTXConfig())

    @property
    def pptx_config(self) -> PPTXConfig:
        """Return the config cast to ``PPTXConfig``."""
        return self.config  # type: ignore[return-value]

    async def load(self, source: str | bytes) -> list[LoadedDocument]:
        """Load a PPTX from a file path or raw bytes.

        In ``PAGE`` mode each slide becomes a separate
        :class:`LoadedDocument`.  In ``SINGLE`` mode (the default) all
        slides are concatenated with slide-number headers.
        """
        from pptx import Presentation  # lazy import — heavy dependency

        if isinstance(source, str):
            try:
                prs = Presentation(source)
            except Exception as exc:
                logger.error("Failed to open PPTX file %s: %s", source, exc)
                raise
            source_name = source
        else:
            self._validate_size(source)
            try:
                prs = Presentation(io.BytesIO(source))
            except Exception as exc:
                logger.error("Failed to parse PPTX from bytes: %s", exc)
                raise
            source_name = "bytes_input"

        slide_count = len(prs.slides)

        if self.pptx_config.mode == LoadMode.PAGE:
            documents: list[LoadedDocument] = []
            for i, slide in enumerate(prs.slides):
                text = self._extract_slide_text(slide)
                if text.strip():
                    documents.append(
                        LoadedDocument(
                            content=text,
                            source=source_name,
                            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            metadata=DocumentMetadata(
                                page_count=slide_count,
                                word_count=self._count_words(text),
                            ),
                            page_number=i + 1,
                        )
                    )
            logger.debug("Loaded %d slides from PPTX %s", len(documents), source_name)
            return documents

        # SINGLE mode — concatenate all slides
        all_text: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_text = self._extract_slide_text(slide)
            if slide_text.strip():
                all_text.append(f"--- Slide {i + 1} ---\n{slide_text}")

        full_text = "\n\n".join(all_text)

        logger.debug("Loaded PPTX %s as single document (%d slides)", source_name, slide_count)
        return [
            LoadedDocument(
                content=full_text,
                source=source_name,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                metadata=DocumentMetadata(
                    page_count=slide_count,
                    word_count=self._count_words(full_text),
                ),
            )
        ]

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _extract_slide_text(self, slide: Any) -> str:
        """Extract all text from a single slide, including tables and notes."""
        parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

            if self.pptx_config.extract_tables and shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))

        if self.pptx_config.extract_notes and slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker Notes] {notes}")

        return "\n".join(parts)

    def supports_source(self, source: str) -> bool:
        return source.lower().endswith(".pptx")


__all__ = [
    "PPTXConfig",
    "PPTXLoader",
]
